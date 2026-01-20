#!/usr/bin/env python3
"""
Convert PREPROCESSING_SLIDES.md to PowerPoint presentation (.pptx)
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import re

def parse_markdown_slides(md_file):
    """Parse markdown file into slide content."""
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Split by slide separators
    slides_raw = re.split(r'^---\s*$', content, flags=re.MULTILINE)
    
    slides = []
    for slide_raw in slides_raw:
        slide_raw = slide_raw.strip()
        if not slide_raw or slide_raw.startswith('# Data Preprocessing'):
            continue
        
        # Extract slide title
        title_match = re.match(r'^# Slide \d+: (.+)$', slide_raw, re.MULTILINE)
        if title_match:
            title = title_match.group(1)
            content = slide_raw[len(title_match.group(0)):].strip()
        else:
            # Check if it's the title slide
            if slide_raw.startswith('# Data Preprocessing'):
                lines = slide_raw.split('\n')
                title = lines[0].replace('# ', '')
                subtitle = lines[1].replace('## ', '') if len(lines) > 1 else ''
                content = '\n'.join(lines[2:]) if len(lines) > 2 else ''
                slides.append({
                    'title': title,
                    'subtitle': subtitle,
                    'content': content,
                    'is_title': True
                })
                continue
            else:
                continue
        
        slides.append({
            'title': title,
            'content': content,
            'is_title': False
        })
    
    return slides

def parse_table(text):
    """Parse markdown table into list of rows."""
    lines = [line.strip() for line in text.split('\n') if line.strip()]
    table_rows = []
    
    for line in lines:
        if '|' in line and not line.startswith('|---'):
            # Remove leading/trailing |
            cells = [cell.strip() for cell in line.split('|')[1:-1]]
            table_rows.append(cells)
    
    return table_rows

def add_text_to_shape(shape, text, is_bold=False, font_size=None):
    """Add text to a text shape with formatting."""
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = text
    if is_bold:
        p.font.bold = True
    if font_size:
        p.font.size = Pt(font_size)
    return p

def create_slide(prs, slide_data):
    """Create a PowerPoint slide from slide data."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Title
    if slide_data.get('is_title'):
        # Title slide
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = slide_data['title']
        title_p = title_frame.paragraphs[0]
        title_p.font.size = Pt(44)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(0, 51, 102)
        title_p.alignment = PP_ALIGN.CENTER
        
        if slide_data.get('subtitle'):
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = slide_data['subtitle']
            subtitle_p = subtitle_frame.paragraphs[0]
            subtitle_p.font.size = Pt(28)
            subtitle_p.font.color.rgb = RGBColor(64, 64, 64)
            subtitle_p.alignment = PP_ALIGN.CENTER
        
        content_y = Inches(4.5)
    else:
        # Regular slide
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = slide_data['title']
        title_p = title_frame.paragraphs[0]
        title_p.font.size = Pt(32)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(0, 51, 102)
        content_y = Inches(1.3)
    
    # Parse content
    content = slide_data.get('content', '')
    
    # Check for tables
    if '|' in content and '---' in content:
        # Extract table
        table_match = re.search(r'\|.*\|.*\n\|.*---.*\n((?:\|.*\|.*\n?)+)', content, re.MULTILINE)
        if table_match:
            table_text = table_match.group(0)
            table_rows = parse_table(table_text)
            
            # Create table
            if len(table_rows) > 0:
                num_cols = len(table_rows[0])
                num_rows = len(table_rows)
                
                table_shape = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), content_y, Inches(9), Inches(2))
                table = table_shape.table
                
                for i, row in enumerate(table_rows):
                    for j, cell_text in enumerate(row):
                        if j < num_cols:
                            cell = table.cell(i, j)
                            cell.text = cell_text
                            # Header row formatting
                            if i == 0:
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
                                for paragraph in cell.text_frame.paragraphs:
                                    paragraph.font.color.rgb = RGBColor(255, 255, 255)
                                    paragraph.font.bold = True
                                    paragraph.font.size = Pt(12)
                            else:
                                for paragraph in cell.text_frame.paragraphs:
                                    paragraph.font.size = Pt(11)
                
                # Remove table from content
                content = content.replace(table_text, '').strip()
                content_y = Inches(4)
    
    # Check for code blocks
    code_blocks = re.findall(r'```(.*?)```', content, re.DOTALL)
    for code_block in code_blocks:
        code_text = code_block.strip()
        code_box = slide.shapes.add_textbox(Inches(0.5), content_y, Inches(9), Inches(1.5))
        code_frame = code_box.text_frame
        code_frame.text = code_text
        code_p = code_frame.paragraphs[0]
        code_p.font.size = Pt(10)
        code_p.font.name = 'Courier New'
        code_p.font.color.rgb = RGBColor(0, 100, 0)
        content = content.replace(f'```{code_block}```', '').strip()
        content_y += Inches(1.8)
    
    # Check for bullet lists
    if any(line.strip().startswith(('-', '*', '1.', '2.', '3.', '4.', '5.')) for line in content.split('\n')):
        lines = [line.strip() for line in content.split('\n') if line.strip()]
        bullet_text = '\n'.join(lines)
        
        text_box = slide.shapes.add_textbox(Inches(0.5), content_y, Inches(9), Inches(4))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        for i, line in enumerate(lines):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            # Remove markdown formatting
            line = re.sub(r'\*\*(.*?)\*\*', r'\1', line)  # Bold
            line = re.sub(r'`(.*?)`', r'\1', line)  # Code
            line = re.sub(r'^[-*]\s*', '', line)  # Bullet
            line = re.sub(r'^\d+\.\s*', '', line)  # Numbered
            
            p.text = line
            p.font.size = Pt(14)
            p.level = 0
            
            if line.startswith('**') or '**' in line:
                p.font.bold = True
    else:
        # Regular text content
        if content.strip():
            # Remove markdown formatting
            content = re.sub(r'\*\*(.*?)\*\*', r'\1', content)  # Bold
            content = re.sub(r'`(.*?)`', r'\1', content)  # Code
            content = re.sub(r'^###\s*', '', content, flags=re.MULTILINE)  # Headers
            content = re.sub(r'^##\s*', '', content, flags=re.MULTILINE)  # Headers
            
            text_box = slide.shapes.add_textbox(Inches(0.5), content_y, Inches(9), Inches(4.5))
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            text_frame.text = content.strip()
            
            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.space_after = Pt(6)

def main():
    """Main function to create PowerPoint presentation."""
    md_file = 'PREPROCESSING_SLIDES.md'
    output_file = 'PREPROCESSING_PRESENTATION.pptx'
    
    print(f"Reading {md_file}...")
    slides_data = parse_markdown_slides(md_file)
    
    print(f"Found {len(slides_data)} slides")
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Create slides
    for i, slide_data in enumerate(slides_data):
        print(f"Creating slide {i+1}: {slide_data.get('title', 'Unknown')}")
        create_slide(prs, slide_data)
    
    # Save
    prs.save(output_file)
    print(f"\n✓ Presentation saved to {output_file}")
    print(f"✓ Created {len(slides_data)} slides")

if __name__ == '__main__':
    main()
